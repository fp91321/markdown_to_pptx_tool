import os
import sys
import tkinter as tk
from tkinter import ttk, scrolledtext, messagebox, filedialog
import re

# 檢查必要套件
try:
    from tkinterdnd2 import TkinterDnD, DND_FILES
    DRAG_DROP_AVAILABLE = True
except ImportError:
    print("警告: tkinterdnd2 未安裝，拖拽功能將不可用")
    print("安裝指令: pip install tkinterdnd2")
    DRAG_DROP_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    print("警告: python-pptx 未安裝，PPTX 轉換功能將不可用")
    print("安裝指令: pip install python-pptx")
    PPTX_AVAILABLE = False

class MarkdownEditorApp(TkinterDnD.Tk if DRAG_DROP_AVAILABLE else tk.Tk):
    def __init__(self):
        try:
            super().__init__()
            self.title("Markdown Editor & PPTX Converter")
            self.geometry("1200x700")
            self.configure(bg="#2d2d2d")
            
            # 變數
            self.current_file = None
            self.font_size = "小"  # 預設字體大小
            
            # 定義字體大小配置
            self.font_sizes = {
                "小": {
                    "h1": 16, "h2": 14, "h3": 12, "h4": 11, 
                    "h5": 10, "h6": 9, "normal": 8, "list": 8, 
                    "code": 7, "quote": 8, "error": 8,
                    "editor": 9, "preview": 8
                },
                "中": {
                    "h1": 20, "h2": 18, "h3": 16, "h4": 15, 
                    "h5": 14, "h6": 13, "normal": 12, "list": 12, 
                    "code": 11, "quote": 12, "error": 12,
                    "editor": 13, "preview": 12
                },
                "大": {
                    "h1": 24, "h2": 22, "h3": 20, "h4": 19, 
                    "h5": 18, "h6": 17, "normal": 16, "list": 16, 
                    "code": 15, "quote": 16, "error": 16,
                    "editor": 17, "preview": 16
                }
            }
            
            self.create_widgets()
            self.setup_bindings()
            print("GUI 初始化完成")
            
        except Exception as e:
            print(f"初始化錯誤: {e}")
            messagebox.showerror("錯誤", f"程式初始化失敗: {e}")
            sys.exit(1)
        
    def create_widgets(self):
        try:
            # 主框架
            main_frame = tk.Frame(self, bg="#2d2d2d")
            main_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
            
            # 標題
            title_label = tk.Label(
                main_frame, 
                text="Markdown Editor & PPTX Converter", 
                font=("Arial", 16, "bold"),
                fg="white", 
                bg="#2d2d2d"
            )
            title_label.pack(pady=(0, 10))
            
            # 工具列
            toolbar = tk.Frame(main_frame, bg="#2d2d2d")
            toolbar.pack(fill=tk.X, pady=(0, 10))
            
            # 文件操作按鈕
            tk.Button(toolbar, text="新建", command=self.new_file).pack(side=tk.LEFT, padx=(0, 5))
            tk.Button(toolbar, text="開啟", command=self.open_file).pack(side=tk.LEFT, padx=(0, 5))
            tk.Button(toolbar, text="儲存(Ctrl+S)", command=self.save_file).pack(side=tk.LEFT, padx=(0, 5))
            
            # 復原/重做按鈕 - ★新增重做功能
            tk.Button(toolbar, text="上一步 (Ctrl+Z)", command=self.undo_action).pack(side=tk.LEFT, padx=(10, 5))
            tk.Button(toolbar, text="下一步 (Ctrl+Y)", command=self.redo_action).pack(side=tk.LEFT, padx=(0, 5))
            
            if PPTX_AVAILABLE:
                tk.Button(toolbar, text="轉換成 PPTX", command=self.convert_to_pptx).pack(side=tk.LEFT, padx=(0, 5))
            
            # 分隔線
            separator_frame = tk.Frame(toolbar, bg="#666666", width=2, height=25)
            separator_frame.pack(side=tk.LEFT, padx=10)
            
            # 字體大小選擇
            tk.Label(toolbar, text="字體大小:", fg="white", bg="#2d2d2d").pack(side=tk.LEFT, padx=(0, 5))
            
            self.font_size_var = tk.StringVar(value=self.font_size)
            font_size_frame = tk.Frame(toolbar, bg="#2d2d2d")
            font_size_frame.pack(side=tk.LEFT, padx=(0, 10))
            
            tk.Radiobutton(font_size_frame, text="小", variable=self.font_size_var, 
                          value="小", command=self.change_font_size,
                          fg="white", bg="#2d2d2d", selectcolor="#2d2d2d").pack(side=tk.LEFT)
            tk.Radiobutton(font_size_frame, text="中", variable=self.font_size_var, 
                          value="中", command=self.change_font_size,
                          fg="white", bg="#2d2d2d", selectcolor="#2d2d2d").pack(side=tk.LEFT)
            tk.Radiobutton(font_size_frame, text="大", variable=self.font_size_var, 
                          value="大", command=self.change_font_size,
                          fg="white", bg="#2d2d2d", selectcolor="#2d2d2d").pack(side=tk.LEFT)
            
            # 預覽模式切換
            self.preview_mode = tk.StringVar(value="rendered")
            
            # ★修改：使用 PanedWindow 實現拖拽調整大小功能
            # 可拖拽的分隔窗格
            self.paned_window = tk.PanedWindow(
                main_frame, 
                orient=tk.HORIZONTAL,
                bg="#2d2d2d",
                sashwidth=8,
                sashrelief=tk.RAISED,
                sashpad=2
            )
            self.paned_window.pack(fill=tk.BOTH, expand=True)
            
            # 左側編輯區框架
            left_frame = tk.LabelFrame(self.paned_window, text="Markdown 編輯器", 
                                     bg="#2d2d2d", fg="white", font=("Arial", 10))
            
            self.editor = scrolledtext.ScrolledText(
                left_frame, 
                wrap=tk.WORD, 
                font=("Consolas", self.font_sizes[self.font_size]["editor"]),
                bg="#1e1e1e", 
                fg="#d4d4d4",
                insertbackground="white",
                undo=True,  # 啟用復原功能
                maxundo=100  # ★增加最大復原步數
            )
            self.editor.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
            
            # 右側預覽區框架
            right_frame = tk.LabelFrame(self.paned_window, text="預覽", 
                                      bg="#2d2d2d", fg="white", font=("Arial", 10))
            
            self.preview = scrolledtext.ScrolledText(
                right_frame, 
                wrap=tk.WORD, 
                font=("Consolas", self.font_sizes[self.font_size]["preview"]),
                bg="white", 
                fg="black",
                state=tk.DISABLED
            )
            self.preview.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
            
            # ★將框架加入 PanedWindow
            self.paned_window.add(left_frame)
            self.paned_window.add(right_frame)
            
            # 狀態列
            self.status_label = tk.Label(
                main_frame, 
                text="準備就緒 - 可拖拽中間分隔線調整視窗大小", 
                font=("Arial", 10),
                fg="white", 
                bg="#2d2d2d",
                anchor=tk.W
            )
            self.status_label.pack(fill=tk.X, pady=(10, 0))
            
            # 作者資訊
            author_label = tk.Label(
                main_frame, 
                text="Editing by Andy Chang, Network Computing Lab., Dept. of Engineering Science, National Cheng Kung University, Taiwan", 
                font=("Arial", 8),
                fg="#888", 
                bg="#2d2d2d"
            )
            author_label.pack(pady=(5, 0))
            
        except Exception as e:
            print(f"建立介面錯誤: {e}")
            raise
        
    def setup_bindings(self):
        try:
            # 編輯器事件綁定
            self.editor.bind('<KeyRelease>', self.on_text_change)
            self.editor.bind('<Button-1>', self.on_text_change)
            
            # 拖拽功能（如果可用）
            if DRAG_DROP_AVAILABLE:
                self.drop_target_register(DND_FILES)
                self.dnd_bind('<<Drop>>', self.drop_file)
            
            # ★增加快捷鍵
            self.bind('<Control-n>', lambda e: self.new_file())
            self.bind('<Control-o>', lambda e: self.open_file())
            self.bind('<Control-s>', lambda e: self.save_file())
            self.bind('<Control-z>', lambda e: self.undo_action())  # Ctrl+Z 復原
            self.bind('<Control-y>', lambda e: self.redo_action())  # Ctrl+Y 重做
            self.bind('<Shift-Control-Z>', lambda e: self.redo_action())  # Shift+Ctrl+Z 重做
            
            # 確保編輯器獲得焦點時快捷鍵也能工作
            self.editor.bind('<Control-z>', lambda e: self.undo_action())
            self.editor.bind('<Control-y>', lambda e: self.redo_action())
            self.editor.bind('<Shift-Control-Z>', lambda e: self.redo_action())
            
        except Exception as e:
            print(f"設定綁定錯誤: {e}")
    
    def change_font_size(self):
        """改變字體大小"""
        try:
            self.font_size = self.font_size_var.get()
            
            # 更新編輯器字體
            editor_font_size = self.font_sizes[self.font_size]["editor"]
            self.editor.config(font=("Consolas", editor_font_size))
            
            # 更新預覽區字體
            preview_font_size = self.font_sizes[self.font_size]["preview"]
            self.preview.config(font=("Consolas", preview_font_size))
            
            # 重新配置標籤並更新預覽
            self._tags_configured = False  # 強制重新配置標籤
            self.update_preview()
            
            self.status_label.config(text=f"字體大小已變更為: {self.font_size}")
            
        except Exception as e:
            print(f"改變字體大小錯誤: {e}")
    
    def undo_action(self):
        """復原操作"""
        try:
            # tkinter 的 Text 元件內建 undo 功能
            self.editor.edit_undo()
            self.status_label.config(text="已復原上一步操作")
            self.update_preview()  # 更新預覽
        except tk.TclError:
            # 沒有可復原的操作
            self.status_label.config(text="沒有可復原的操作")
        except Exception as e:
            print(f"復原操作錯誤: {e}")
    
    # ★新增重做功能
    def redo_action(self):
        """重做操作"""
        try:
            # tkinter 的 Text 元件內建 redo 功能
            self.editor.edit_redo()
            self.status_label.config(text="已重做操作")
            self.update_preview()  # 更新預覽
        except tk.TclError:
            # 沒有可重做的操作
            self.status_label.config(text="沒有可重做的操作")
        except Exception as e:
            print(f"重做操作錯誤: {e}")
        
    def on_text_change(self, event=None):
        """當文字改變時更新預覽"""
        try:
            self.after_idle(self.update_preview)
        except Exception as e:
            print(f"文字改變處理錯誤: {e}")
        
    def update_preview(self):
        """更新預覽區域 - 支援字體大小"""
        try:
            content = self.editor.get("1.0", tk.END)
            
            self.preview.config(state=tk.NORMAL)
            self.preview.delete("1.0", tk.END)
            
            if self.preview_mode.get() == "rendered":
                # 配置文字標籤
                if not hasattr(self, '_tags_configured') or not self._tags_configured:
                    self.configure_text_tags()
                    self._tags_configured = True
                
                # 獲取格式化數據
                formatted_data = self.parse_markdown_to_text(content)
                
                # 逐段插入帶有標籤的文字
                for text, tag in formatted_data:
                    self.preview.insert(tk.END, text, tag)
            else:
                # 原始 Markdown 預覽
                self.preview.insert("1.0", content)
                
            self.preview.config(state=tk.DISABLED)
            
        except Exception as e:
            print(f"更新預覽錯誤: {e}")
            self.preview.config(state=tk.NORMAL)
            self.preview.delete("1.0", tk.END)
            self.preview.insert("1.0", f"預覽錯誤: {str(e)}", "error")
            self.preview.config(state=tk.DISABLED)
        
    def parse_markdown_to_text(self, markdown_text):
        """直接解析 Markdown 並保持格式，返回格式化信息"""
        try:
            lines = markdown_text.splitlines()
            formatted_data = []  # 存儲 (text, tag) 的列表
            
            for line in lines:
                if not line.strip():
                    formatted_data.append(("", "normal"))
                    continue
                    
                # 標題處理 - 不同層級使用不同字體大小
                if line.startswith("# "):
                    text = "\n═══ " + line[2:] + " ═══\n\n"
                    formatted_data.append((text, "h1"))
                elif line.startswith("## "):
                    text = "\n▶ " + line[3:] + "\n"
                    formatted_data.append((text, "h2"))
                elif line.startswith("### "):
                    text = "\n● " + line[4:] + "\n"
                    formatted_data.append((text, "h3"))
                elif line.startswith("#### "):
                    text = "◆ " + line[5:] + "\n"
                    formatted_data.append((text, "h4"))
                elif line.startswith("##### "):
                    text = "◇ " + line[6:] + "\n"
                    formatted_data.append((text, "h5"))
                elif line.startswith("###### "):
                    text = "○ " + line[7:] + "\n"
                    formatted_data.append((text, "h6"))
                
                # 列表處理 - 保持原始縮排
                elif line.lstrip().startswith(("- ", "* ", "+ ")):
                    leading_spaces = len(line) - len(line.lstrip())
                    indent = " " * leading_spaces
                    content = line.lstrip()[2:]
                    text = f"{indent}• {content}\n"
                    formatted_data.append((text, "list"))
                    
                # 數字列表處理
                elif re.match(r'^\s*\d+\.\s', line):
                    leading_spaces = len(line) - len(line.lstrip())
                    indent = " " * leading_spaces
                    match = re.match(r'^(\s*)(\d+)\.\s(.*)$', line)
                    if match:
                        number = match.group(2)
                        content = match.group(3)
                        text = f"{indent}{number}. {content}\n"
                        formatted_data.append((text, "list"))
                    else:
                        formatted_data.append((line + "\n", "normal"))
                        
                # ★修正：代碼塊處理（補上冒號）
                elif line.strip().startswith("```"):
                    if line.strip() == "```":
                        text = "─" * 50 + "\n"
                        formatted_data.append((text, "code"))
                    else:
                        lang = line.strip()[3:]
                        text = f"─── {lang} ───\n"
                        formatted_data.append((text, "code"))
                        
                # 引用處理
                elif line.lstrip().startswith("> "):
                    leading_spaces = len(line) - len(line.lstrip())
                    indent = " " * leading_spaces
                    content = line.lstrip()[2:]
                    text = f"{indent}│ {content}\n"
                    formatted_data.append((text, "quote"))
                    
                # 普通段落
                else:
                    formatted_data.append((line + "\n", "normal"))
                    
            return formatted_data
            
        except Exception as e:
            return [("解析錯誤: " + str(e), "error")]

    def configure_text_tags(self):
        """配置文字標籤的字體和顏色"""
        # 基本字體設定
        base_font_family = "Consolas"
        sizes = self.font_sizes[self.font_size]
        
        # 配置不同標籤的樣式
        self.preview.tag_configure("h1", font=(base_font_family, sizes["h1"], "bold"), foreground="#000080")
        self.preview.tag_configure("h2", font=(base_font_family, sizes["h2"], "bold"), foreground="#000080")
        self.preview.tag_configure("h3", font=(base_font_family, sizes["h3"], "bold"), foreground="#000080")
        self.preview.tag_configure("h4", font=(base_font_family, sizes["h4"], "bold"), foreground="#000080")
        self.preview.tag_configure("h5", font=(base_font_family, sizes["h5"], "bold"), foreground="#000080")
        self.preview.tag_configure("h6", font=(base_font_family, sizes["h6"], "bold"), foreground="#000080")
        self.preview.tag_configure("normal", font=(base_font_family, sizes["normal"]), foreground="#000000")
        self.preview.tag_configure("list", font=(base_font_family, sizes["list"]), foreground="#333333")
        self.preview.tag_configure("code", font=(base_font_family, sizes["code"]), foreground="#800080", background="#f0f0f0")
        self.preview.tag_configure("quote", font=(base_font_family, sizes["quote"], "italic"), foreground="#666666")
        self.preview.tag_configure("error", font=(base_font_family, sizes["error"]), foreground="#ff0000")

    def new_file(self):
        """新建文件"""
        try:
            if messagebox.askyesno("新建", "確定要新建文件嗎？未儲存的內容將丟失。"):
                self.editor.delete("1.0", tk.END)
                self.current_file = None
                self.status_label.config(text="新建文件")
        except Exception as e:
            print(f"新建文件錯誤: {e}")
            
    def open_file(self):
        """開啟文件"""
        try:
            file_path = filedialog.askopenfilename(
                title="開啟 Markdown 文件",
                filetypes=[("Markdown files", "*.md"), ("Text files", "*.txt"), ("All files", "*.*")]
            )
            if file_path:
                self.load_file(file_path)
        except Exception as e:
            print(f"開啟文件錯誤: {e}")
            messagebox.showerror("錯誤", f"開啟文件失敗: {e}")
            
    def load_file(self, file_path):
        """載入文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            self.editor.delete("1.0", tk.END)
            self.editor.insert("1.0", content)
            self.current_file = file_path
            self.status_label.config(text=f"已開啟: {os.path.basename(file_path)}")
            self.update_preview()
        except Exception as e:
            print(f"載入文件錯誤: {e}")
            messagebox.showerror("錯誤", f"無法載入文件: {str(e)}")
            
    def save_file(self):
        """儲存文件"""
        try:
            if self.current_file:
                self.save_to_file(self.current_file)
            else:
                self.save_as_file()
        except Exception as e:
            print(f"儲存文件錯誤: {e}")
            
    def save_as_file(self):
        """另存新檔"""
        try:
            file_path = filedialog.asksaveasfilename(
                title="儲存 Markdown 文件",
                defaultextension=".md",
                filetypes=[("Markdown files", "*.md"), ("Text files", "*.txt"), ("All files", "*.*")]
            )
            if file_path:
                self.save_to_file(file_path)
                self.current_file = file_path
        except Exception as e:
            print(f"另存新檔錯誤: {e}")
            
    def save_to_file(self, file_path):
        """儲存到指定文件"""
        try:
            content = self.editor.get("1.0", tk.END)
            with open(file_path, 'w', encoding='utf-8') as file:
                file.write(content)
            self.status_label.config(text=f"已儲存: {os.path.basename(file_path)}")
        except Exception as e:
            print(f"儲存到文件錯誤: {e}")
            messagebox.showerror("錯誤", f"無法儲存文件: {str(e)}")
            
    def convert_to_pptx(self):
        """轉換成 PPTX"""
        if not PPTX_AVAILABLE:
            messagebox.showerror("錯誤", "python-pptx 套件未安裝")
            return
            
        try:
            content = self.editor.get("1.0", tk.END).strip()
            if not content:
                messagebox.showwarning("警告", "請先輸入 Markdown 內容")
                return
                
            file_path = filedialog.asksaveasfilename(
                title="儲存 PPTX 文件",
                defaultextension=".pptx",
                filetypes=[("PowerPoint files", "*.pptx"), ("All files", "*.*")]
            )
            
            if file_path:
                self.markdown_to_pptx(content, file_path)
                self.status_label.config(text=f"PPTX 轉換完成: {os.path.basename(file_path)}")
                messagebox.showinfo("成功", f"已成功轉換為 {os.path.basename(file_path)}")
                
        except Exception as e:
            error_msg = f"轉換失敗: {str(e)}"
            print(error_msg)
            self.status_label.config(text=error_msg)
            messagebox.showerror("錯誤", error_msg)
            
    def drop_file(self, event):
        """處理拖拽文件"""
        try:
            file_path = event.data.strip('{}')
            if file_path.lower().endswith(('.md', '.txt')):
                self.load_file(file_path)
            else:
                messagebox.showwarning("警告", "請拖拽 .md 或 .txt 文件")
        except Exception as e:
            print(f"拖拽處理錯誤: {e}")
            
    def markdown_to_pptx(self, md_text, output_file):
        """Markdown 轉 PPTX"""
        try:
            presentation = Presentation()
            slide = None
            
            lines = md_text.splitlines()
            for line in lines:
                line = line.rstrip()
                if line.startswith("# "):
                    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
                    slide.shapes.title.text = line[2:]
                elif line.startswith("## "):
                    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
                    slide.shapes.title.text = line[3:]
                elif line.startswith("### "):
                    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
                    slide.shapes.title.text = line[4:]
                elif line.strip().startswith(("-", "*")):
                    if not slide:
                        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
                        slide.shapes.title.text = "Bullet Points"
                    
                    if len(slide.shapes.placeholders) > 1 and slide.shapes.placeholders[1].has_text_frame:
                        text_frame = slide.shapes.placeholders[1].text_frame
                    else:
                        text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
                        text_frame = text_box.text_frame
                    
                    current_level = (len(line) - len(line.lstrip())) // 4
                    content = line.strip("-* ").strip()
                    
                    p = text_frame.add_paragraph()
                    p.text = content
                    p.level = min(current_level, 4)
                    
            presentation.save(output_file)
            
        except Exception as e:
            raise Exception(f"PPTX 轉換錯誤: {str(e)}")

def main():
    """主函數，包含錯誤處理"""
    try:
        print("正在啟動 Markdown 編輯器...")
        app = MarkdownEditorApp()
        print("啟動成功，開始主迴圈...")
        app.mainloop()
    except Exception as e:
        print(f"程式執行錯誤: {e}")
        import traceback
        traceback.print_exc()
        input("按 Enter 鍵退出...")

if __name__ == "__main__":
    main()
